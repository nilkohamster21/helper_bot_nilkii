from pptx import Presentation
from pptx.util import Cm, Pt
import json

def generate_presentation(text_blocks, image_paths, template_config_path, title_output):
    # загрузка шаблона конфигурации
    with open(template_config_path, 'r') as f:
        template_config = json.load(f)

    # создание новой презентации
    prs = Presentation()

    text_index = 0  # индекс текущего текста
    image_index = 0  # индекс текущего изображения

    # слайд 1 титульный
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # стандартный титульный шаблон
    title = slide1.shapes.title

    # установка заголовка
    if text_index < len(text_blocks):
        title.text = text_blocks[text_index]
        title.text_frame.paragraphs[0].font.size = Pt(40) # размер шрифта
        title.text_frame.paragraphs[0].font.bold = False  # убираю жирность
        text_index += 1

    # слайд 2 пустой слайд
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # добавление текстов на слайд 2
    for tb in template_config['slide2']['text_blocks']:
        if text_index < len(text_blocks):
            textbox = slide2.shapes.add_textbox(Cm(tb['x']), Cm(tb['y']), Cm(tb['width']), Cm(tb['height']))
            frame = textbox.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            p.text = text_blocks[text_index]
            p.font.size = Pt(24)
            text_index += 1

    # добавление картинок на слайд 2
    for img in template_config['slide2']['images']:
        if image_index < len(image_paths):
            pic_path = image_paths[image_index]
            try:
                slide2.shapes.add_picture(pic_path, Cm(img['x']), Cm(img['y']),
                                          width=Cm(img['max_width']), height=None)
            except Exception as e:
                print(f"Ошибка с картинкой: {pic_path}, {e}")
            image_index += 1

    # слайд 3 пустой слайд
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # добавление текстов на слайд 3
    for tb in template_config['slide3']['text_blocks']:
        if text_index < len(text_blocks):
            textbox = slide3.shapes.add_textbox(Cm(tb['x']), Cm(tb['y']), Cm(tb['width']), Cm(tb['height']))
            frame = textbox.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            p.text = text_blocks[text_index]
            p.font.size = Pt(24)
            text_index += 1

    # добавление картинок на слайд 3
    for img in template_config['slide3']['images']:
        if image_index < len(image_paths):
            pic_path = image_paths[image_index]
            try:
                slide3.shapes.add_picture(pic_path, Cm(img['x']), Cm(img['y']),
                                          width=Cm(img['max_width']), height=None)
            except Exception as e:
                print(f"Ошибка с картинкой: {pic_path}, {e}")
            image_index += 1

    # сохранение презентации
    prs.save(title_output)
    print(f"Презентация сохранена как {title_output}")
    return title_output
